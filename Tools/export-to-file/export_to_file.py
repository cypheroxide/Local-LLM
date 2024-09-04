"""
LLM Output Converter
---------------------
A Python script to convert LLM outputs into common document formats:
- Word (.docx)
- PowerPoint (.pptx)
- Excel (.xlsx)

Usage:
    - Save as Word Document: `save_as_word(text, "output.docx")`
    - Save as PowerPoint Presentation: `save_as_ppt(text, "output.pptx")`
    - Save as Excel Spreadsheet: `save_as_excel(data, "output.xlsx")`

Requirements:
    - python-docx
    - python-pptx
    - pandas
    - openpyxl

Author: Cypher Oxide
GitHub: https://github.com/cypheroxide

License: MIT License
"""

from docx import Document
from pptx import Presentation
import pandas as pd

def save_as_word(text, filename):
    doc = Document()
    doc.add_paragraph(text)
    doc.save(filename)

def save_as_ppt(text, filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Chat Output"
    content.text = text
    
    prs.save(filename)

def save_as_excel(data, filename):
    df = pd.DataFrame(data, columns=['Chat Output'])
    df.to_excel(filename, index=False)

# Example usage
chat_output = "This is the LLM output text."
data = [[chat_output]]

# Save as Word
save_as_word(chat_output, "output.docx")

# Save as PowerPoint
save_as_ppt(chat_output, "output.pptx")

# Save as Excel
save_as_excel(data, "output.xlsx")
